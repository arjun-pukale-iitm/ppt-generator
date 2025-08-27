from pptx import Presentation
from pptx.util import Inches
import io

def inspect_template_bytes(bytes_io):
    """Inspect uploaded PPTX template: layouts + images."""
    prs = Presentation(bytes_io)
    layouts = []
    images = []

    for i, layout in enumerate(prs.slide_layouts):
        layouts.append({
            "index": i,
            "placeholder_count": len(layout.placeholders)
        })

    for s in prs.slides:
        for shape in s.shapes:
            if hasattr(shape, "image"):
                images.append({"name": getattr(shape, "name", f"img_{len(images)}")})

    return {"layouts": layouts, "images": images}


def build_presentation_from_plan(template_bytes, plan_json):
    """Generate PPTX from template + JSON plan, respecting formatting."""
    prs = Presentation(template_bytes)

    # Remove all pre-existing slides from template (keep layouts only)
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # Cache template images (if any)
    template_images = {}
    for layout in prs.slide_layouts:
        for shape in layout.shapes:
            if hasattr(shape, "image"):
                name = getattr(shape, "name", f"img_{len(template_images)}")
                template_images[name] = shape.image.blob

    # Build new slides from the plan
    for slide_def in plan_json.get("slides", []):
        layout_idx = slide_def.get("layout_index", 0)
        layout = prs.slide_layouts[layout_idx] if layout_idx < len(prs.slide_layouts) else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        # Fill placeholders with content
        for ph in slide.placeholders:
            try:
                # Title placeholder
                if ph.placeholder_format.type == 1:  
                    ph.text = slide_def.get("title", "")

                # Content placeholder
                elif ph.placeholder_format.type == 2:
                    tf = ph.text_frame
                    tf.clear()
                    bullets = slide_def.get("bullets", [])
                    if bullets:
                        tf.paragraphs[0].text = bullets[0]
                        for b in bullets[1:]:
                            tf.add_paragraph().text = b
            except:
                continue

        # Optional image insertion from template hints
        img_tag = slide_def.get("image_from_template_hint")
        if img_tag and template_images:
            blob = next((img for name, img in template_images.items() if img_tag.lower() in name.lower()), None)
            if not blob:
                blob = next(iter(template_images.values()))
            slide.shapes.add_picture(io.BytesIO(blob), Inches(5), Inches(1.5), width=Inches(4))

        # Add speaker notes
        if slide_def.get("notes"):
            slide.notes_slide.notes_text_frame.text = slide_def["notes"]

    # Save final presentation
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out
